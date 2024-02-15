#python-pptx == 0.6.23

from pptx import Presentation
from pptx.util import Inches, Pt

import tkinter as tk
from tkinter import filedialog
from tkinter import StringVar

import os




def add_slide_page(pptx_file_path, first_count, localte_per_from_left, localte_per_from_top, save_file_path, size_of_font,flag):
    presentation = Presentation(pptx_file_path)

    rate = float(Inches(1))
    
    left = Inches(localte_per_from_left*(presentation.slide_width)/(100*rate))
    top = Inches(localte_per_from_top*(presentation.slide_height)/(100*rate))
    width = Inches((100 - localte_per_from_left)*(presentation.slide_width)/(100*rate))
    height = Inches((100 - localte_per_from_top)*(presentation.slide_height)/(100*rate))
    
    n = first_count
    all_sliders_num = len(presentation.slides) + n - 1

    for slide in presentation.slides:
        textbox = slide.shapes.add_textbox(left, top, width, height)
        if flag == 1:
            textbox.text_frame.text = f"{n}/{all_sliders_num}"
        else:
            textbox.text_frame.text = f"{n}"
        n += 1
        textbox.text_frame.paragraphs[0].font.size = Pt(size_of_font)

    presentation.save(save_file_path)


def main_run():
    pptx_file_path = filedialog.askopenfilename(filetypes=[("テキストファイル", "*.pptx")])
    directory, filename = os.path.split(pptx_file_path)
    
    if entry_first_num.get() == "":
        first_count = int(1)
    else:
        first_count = int(entry_first_num.get())
    
    if entry_left_local.get() == "":
        localte_per_from_left = float(90)
    else:
        localte_per_from_left = float(entry_left_local.get())
    
    if entry_top_local.get() == "":
        localte_per_from_top = float(90)
    else:
        localte_per_from_top = float(entry_top_local.get())

    flag = int(radioValue.get())

    if entry_font_size.get() == "":
        font_size = float(18)
    else:
        font_size = float(entry_font_size.get())
    
    try:
        file_name, extension = os.path.splitext(filename)
        save_file = f"{file_name}_AddPageNumber{extension}"
        save_file_path = os.path.join(directory, save_file)
        
        i = 1
        while os.path.exists(save_file_path):
            save_file = f"{file_name}_AddPageNumber({i}){extension}"
            save_file_path = os.path.join(directory, save_file)
            i += 1
        
        add_slide_page(
            pptx_file_path, first_count, localte_per_from_left, localte_per_from_top, save_file_path, font_size, flag
            )
        message = f"正常に動作しました．\n{save_file_path}\nとして保存されました．"
    
    except Exception as e:
        error_message = str(e)
        if f"Package not found at '{pptx_file_path}'" in error_message:
            message = "対象のパワーポイントファイルが見つかりませんでした．\n対象のパワーポイントファイルを開いている場合は閉じてください．"
        else:
            message = "予期せぬエラーが発生しました．"
            print(e)
    run_label.config(text=message)
    

def run():
    if entry_first_num.get() == "":
        print(1)
    print(entry_first_num.get())



if __name__ == "__main__":
    root = tk.Tk()
    root.title("Add page number")

    lavel_first_num = tk.Label(root, text="\n開始する数字を入力してください")
    lavel_first_num.grid()
    entry_first_num = tk.Entry(root)
    entry_first_num.grid()

    
    lavel_left_local = tk.Label(root, text="\n左端からの位置をパーセンテージで入力してください")
    lavel_left_local.grid()
    entry_left_local = tk.Entry(root)
    entry_left_local.grid()

    lavel_top_local = tk.Label(root, text="\n上端からの位置をパーセンテージで入力してください")
    lavel_top_local.grid()
    entry_top_local = tk.Entry(root)
    entry_top_local.grid()

    lavel_font_size = tk.Label(root, text="\nフォントのサイズを入力してください")
    lavel_font_size.grid()
    entry_font_size = tk.Entry(root)
    entry_font_size.grid()
    
    radioValue = StringVar()
    radioValue.set(0)
    radio_button_with = tk.Radiobutton(root, text="分母付き", variable=radioValue, value=1)
    radio_button_with.grid()
    radio_button_out = tk.Radiobutton(root, text="分母無し", variable=radioValue, value=0)
    radio_button_out.grid()

    run_button = tk.Button(root, text = "実行", command = main_run)
    run_button.grid()

    run_label = tk.Label(root, text="")
    run_label.grid()

    root.mainloop()
